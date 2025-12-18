import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# Project paths
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
APP_PY = os.path.join(PROJECT_ROOT, 'app.py')
TEMPLATES_DIR = os.path.join(PROJECT_ROOT, 'templates')


def collect_routes(app_py_path):
    routes = []
    try:
        with open(app_py_path, 'r', encoding='utf-8') as f:
            text = f.read()
        for m in re.finditer(r"@app\.route\(\s*['\"]([^'\"]+)['\"]", text):
            routes.append(m.group(1))
    except Exception:
        pass
    return sorted(set(routes))


def list_templates(tmpl_dir):
    files = []
    if not os.path.isdir(tmpl_dir):
        return files
    for root, dirs, filenames in os.walk(tmpl_dir):
        for fn in filenames:
            if fn.endswith('.html'):
                rel = os.path.relpath(os.path.join(root, fn), tmpl_dir)
                files.append(rel.replace('\\', '/'))
    return sorted(files)


def extract_db_tables(app_py_path):
    tables = []
    try:
        with open(app_py_path, 'r', encoding='utf-8') as f:
            text = f.read()
        for m in re.finditer(r"CREATE TABLE IF NOT EXISTS\s+([a-zA-Z0-9_]+)", text, re.IGNORECASE):
            tables.append(m.group(1))
    except Exception:
        pass
    return sorted(set(tables))


# Presentation utilities
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BRAND = RGBColor(10, 30, 60)
MUTED = RGBColor(100, 100, 100)


def title_slide(title, subtitle, student_name='[Student Name]', student_number='[Student #]'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRAND

    sb = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(11), Inches(0.8))
    sp = sb.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = MUTED

    info = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(11), Inches(0.6))
    ip = info.text_frame.paragraphs[0]
    ip.text = f"Student: {student_name}    |    Student #: {student_number}"
    ip.font.size = Pt(12)
    ip.font.color.rgb = MUTED


def bullets_slide(heading, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    h = head.text_frame.paragraphs[0]
    h.text = heading
    h.font.size = Pt(28)
    h.font.bold = True
    h.font.color.rgb = BRAND

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.level = 0


def build_presentation(student_name='[Student Name]', student_number='[Student #]'):
    # Gather some repo info to use where helpful
    routes = collect_routes(APP_PY)
    templates = list_templates(TEMPLATES_DIR)
    tables = extract_db_tables(APP_PY)

    # Slide 0: Title
    title_slide('OCONNECT', 'Media sharing platform with dark-first UI and per-user variants', student_name, student_number)

    # Slides 1-2: Problem & scalability issues
    bullets_slide('Problem Statement', [
        'Users need a reliable, low-latency platform for sharing media and interactions',
        'Requirements: availability, low latency for media, consistent notifications and comments',
        'Scalability challenges: media storage growth, DB contention, and real-time features'
    ])
    bullets_slide('Scalability Issues Identified', [
        'Local file storage causes operational scaling limits and backup complexity',
        'Synchronous processing of uploads and notifications creates throughput bottlenecks',
        'Stateful real-time components (if used) complicate horizontal scaling'
    ])

    # Slides 3-6: Technical solution overview (4 slides)
    bullets_slide('Solution: Architecture Overview', [
        'Stateless Flask app instances behind a load balancer',
        'Offload media to object storage (Blob/S3) and serve via CDN',
        'Background workers for processing (thumbnails, transcoding, notifications)'
    ])
    bullets_slide('Solution: Data & Storage', [
        f"Primary DB (SQLite in project) for prototyping; recommend managed SQL for production",
        'Use object storage for uploads and CDN for global distribution',
        'Use caching (Redis) to reduce DB read pressure for hot content'
    ])
    bullets_slide('Solution: UI / UX & Features', [
        'Per-user template variants (templates/v2) and dark-only default',
        'Numeric rating (1-10) and share counts implemented in server and templates',
        'Nested comments and comment likes implemented server-side and rendered in templates'
    ])
    bullets_slide('Solution: Operations & CI', [
        'Automate DB migrations and backups before deploying upgrades',
        'CI pipeline: run linters, tests, and package app; deploy to container registry',
        'Instrument with Application Insights / Prometheus + Grafana for metrics'
    ])

    # Slides 7-8: Advanced features
    bullets_slide('Advanced Features (1)', [
        'Per-user UI variants: two visual variants supported via `users.ui_variant`',
        'Share count tracking with server-side increment and client updates',
        'Creator pages showing creator posts grid and share controls'
    ])
    bullets_slide('Advanced Features (2)', [
        'Nested, threaded comments with reply and like toggles',
        'Rating upsert (1-10) per user per media item',
        'Optional real-time messaging was present; project can run without Socket.IO for scaling'
    ])

    # Slides 9-10: Limitations & evaluation of ability to scale
    bullets_slide('Limitations (1)', [
        'SQLite is suitable for prototype but not for horizontal scale; upgrade needed',
        'Local static file storage limits scaling and regional redundancy',
        'Embedded Jinja in JS causes maintainability and static-analysis warnings'
    ])
    bullets_slide('Scalability Evaluation (2)', [
        'Stateless web tier + blob storage + CDN decouples traffic from DB load',
        'Background workers remove synchronous bottlenecks for media processing',
        'Database scaling requires managed SQL or NoSQL with partitioning for high traffic'
    ])

    # Slide 11: Demo (5-minute video) and functionality to show
    bullets_slide('Demo Functionality (5-minute video)', [
        'Record: signup -> upload media -> browse feed -> rate -> comment -> share',
        'Show creator upload grid and comment reply/like flows',
        'Show deployment snapshot: services used (App Service / Container Apps, Blob Storage, CDN)'
    ])

    # Slide 12: Conclusions
    bullets_slide('Conclusions & Next Steps', [
        'Migrate storage to blob storage and replace SQLite with managed SQL for production',
        'Add CI/CD and automated tests to maintain quality',
        'Replace risky template-in-JS patterns and improve UX (AJAX rating, modal share)'
    ])

    # Slide 13: References
    bullets_slide('References', [
        'python-pptx: https://python-pptx.readthedocs.io/',
        'Azure Architecture Center: https://learn.microsoft.com/azure/architecture',
        'Flask docs: https://flask.palletsprojects.com/'
    ])


def main():
    # Allow optional environment variables for student info
    sname = os.environ.get('OCONNECT_STUDENT_NAME', '[Student Name]')
    snum = os.environ.get('OCONNECT_STUDENT_NUMBER', '[Student #]')
    build_presentation(sname, snum)
    out = os.path.join(PROJECT_ROOT, 'OCONNECT_presentation.pptx')
    prs.save(out)
    print('Presentation written to', out)

if __name__ == '__main__':
    main()

